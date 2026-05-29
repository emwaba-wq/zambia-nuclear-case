"""Side-by-side diff of hand-edited vs script-regenerated deck."""
from pptx import Presentation


def fingerprint(path):
    prs = Presentation(path)
    out = []
    for i, sl in enumerate(prs.slides, 1):
        shapes = list(sl.shapes)
        text_blocks = []
        for sh in shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    if p.text.strip():
                        text_blocks.append(p.text.strip())
        pic_n = sum(1 for sh in shapes if sh.shape_type == 13)
        chart_n = sum(1 for sh in shapes if sh.has_chart)
        out.append({
            "slide": i,
            "shapes": len(shapes),
            "pics": pic_n,
            "charts": chart_n,
            "texts": text_blocks,
        })
    return out


A = fingerprint("The_Case_for_Nuclear_Power_in_Zambia.HAND_EDITED.pptx")
B = fingerprint("The_Case_for_Nuclear_Power_in_Zambia.REGEN.pptx")

print(f"Slides — hand: {len(A)}  regen: {len(B)}\n")
for i in range(max(len(A), len(B))):
    a = A[i] if i < len(A) else None
    b = B[i] if i < len(B) else None
    if a is None:
        print(f"SLIDE {i+1}: missing in HAND")
        continue
    if b is None:
        print(f"SLIDE {i+1}: missing in REGEN")
        continue
    print(f"=== SLIDE {a['slide']} ===")
    print(f"  shapes  hand={a['shapes']:3d}  regen={b['shapes']:3d}  "
          f"diff={b['shapes']-a['shapes']:+d}")
    print(f"  pics    hand={a['pics']:3d}  regen={b['pics']:3d}")
    print(f"  charts  hand={a['charts']:3d}  regen={b['charts']:3d}")
    # Text diffs
    hand_set = set(a["texts"])
    regen_set = set(b["texts"])
    only_hand = hand_set - regen_set
    only_regen = regen_set - hand_set
    if only_hand:
        print(f"  -- text only in HAND ({len(only_hand)}):")
        for s in sorted(only_hand)[:10]:
            print(f"     - {s[:100]}")
    if only_regen:
        print(f"  -- text only in REGEN ({len(only_regen)}):")
        for s in sorted(only_regen)[:10]:
            print(f"     + {s[:100]}")
    if not only_hand and not only_regen:
        print(f"  text: IDENTICAL ({len(hand_set)} blocks)")
    print()
