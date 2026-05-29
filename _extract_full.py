"""Extract media files + full text from current .pptx."""
import os
import zipfile
import shutil
from pptx import Presentation

P = "The_Case_for_Nuclear_Power_in_Zambia.pptx"
MEDIA = "media"

# 1) Extract embedded media from the pptx ZIP
os.makedirs(MEDIA, exist_ok=True)
with zipfile.ZipFile(P) as z:
    for name in z.namelist():
        if name.startswith("ppt/media/") and not name.endswith("/"):
            target = os.path.join(MEDIA, os.path.basename(name))
            if not os.path.basename(name):
                continue
            with z.open(name) as src, open(target, "wb") as dst:
                shutil.copyfileobj(src, dst)
            print(f"extracted: {target}")

# 2) Print full text on key slides (1, 4, 9, 10)
prs = Presentation(P)
for idx in [1, 4, 9, 10]:
    slide = prs.slides[idx - 1]
    print(f"\n--- FULL TEXT SLIDE {idx} ---")
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                if p.text.strip():
                    print(repr(p.text))

# 3) Print picture rel info for each slide so we know which file is where
print("\n--- PICTURES ---")
for i, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE
            try:
                blob_name = sh.image.filename
            except Exception:
                blob_name = "?"
            print(f"slide {i}: L={sh.left/914400:.2f} T={sh.top/914400:.2f} "
                  f"W={sh.width/914400:.2f} H={sh.height/914400:.2f} "
                  f"src={blob_name}")
