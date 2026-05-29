"""
build_presentation.py
Reconstructs The_Case_for_Nuclear_Power_in_Zambia.pptx
to mirror the hand-edited 10-slide deck (16:9, 10 x 5.625 in).

Icons referenced from ./media/ — extracted from the source deck.
If a media file is missing, that picture is skipped gracefully.

Run:
  pip install -r requirements.txt
  python build_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION,
                             XL_LABEL_POSITION)

OUTPUT = "The_Case_for_Nuclear_Power_in_Zambia.pptx"
MEDIA = "media"

# Palette (sampled from current deck)
NAVY_DEEP = RGBColor(0x0B, 0x1D, 0x2E)
NAVY = RGBColor(0x1A, 0x3A, 0x5C)
NAVY_PANEL = RGBColor(0x14, 0x2C, 0x45)
TEAL = RGBColor(0x0D, 0x73, 0x77)
TEAL_GREEN = RGBColor(0x14, 0xA0, 0x85)
TEAL_LIGHT = RGBColor(0xC8, 0xE8, 0xE4)
TEAL_MINT = RGBColor(0xA8, 0xF0, 0xE8)
TEAL_PALE = RGBColor(0xE8, 0xF8, 0xF5)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
CRISIS = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
BG_BLUE = RGBColor(0xF0, 0xF4, 0xF8)
SLATE = RGBColor(0x8B, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_GREY = RGBColor(0x66, 0x66, 0x66)
PINK_LIGHT = RGBColor(0xFF, 0xB3, 0xB3)
PINK_PALE = RGBColor(0xFF, 0xE0, 0xE0)
BLANK_LAYOUT = 6


# ---------- helpers ----------
def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(l), Inches(t),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def textbox(slide, l, t, w, h, runs, *,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t),
                                  Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Inches(0))
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in runs:
        text = r["text"]
        if text.startswith("\n"):
            p = tf.add_paragraph()
            p.alignment = r.get("align", align)
            p.line_spacing = r.get("line_spacing", line_spacing)
            text = text[1:]
        run = p.add_run()
        run.text = text
        run.font.name = r.get("font", "Calibri")
        run.font.size = Pt(r["size"])
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", NAVY_DEEP)
    return tb


def txt(slide, l, t, w, h, text, *,
        size=12, bold=False, italic=False, color=NAVY_DEEP,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        font="Calibri", line_spacing=1.2):
    return textbox(slide, l, t, w, h,
                   [{"text": text, "size": size, "bold": bold,
                     "italic": italic, "color": color, "font": font}],
                   align=align, anchor=anchor, line_spacing=line_spacing)


def picture(slide, fname, l, t, w, h):
    path = os.path.join(MEDIA, fname)
    if not os.path.isfile(path):
        return None
    return slide.shapes.add_picture(path,
                                    Inches(l), Inches(t),
                                    Inches(w), Inches(h))


def header_bar(slide, fill_color, eyebrow, eyebrow_color, title,
               title_color=WHITE):
    rect(slide, 0, 0, 10, 1.1, fill_color)
    txt(slide, 0.5, 0.08, 9.0, 0.35,
        eyebrow, size=9.5, bold=True, color=eyebrow_color)
    txt(slide, 0.5, 0.42, 9.0, 0.55,
        title, size=26, bold=True, color=title_color,
        font="Georgia")


def color_points(series, colors):
    for i, p in enumerate(series.points):
        if i >= len(colors):
            break
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = colors[i]
        p.format.line.fill.background()


# ---------- SLIDE 1 — Title ----------
def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, NAVY_DEEP)

    rect(s, 0, 0, 0.12, 5.625, GOLD)
    rect(s, 5.5, -1.5, 6.5, 6.5, NAVY)
    picture(s, "image-1-1.png", 7.8, 1.0, 1.6, 1.6)

    txt(s, 0.35, 0.9, 7.0, 0.4,
        "NATIONAL ENERGY STRATEGY",
        size=11, bold=True, color=GOLD)
    textbox(
        s, 0.35, 1.4, 8.5, 2.0,
        [
            {"text": "Zambia Needs", "size": 44, "bold": True,
             "color": WHITE, "font": "Georgia"},
            {"text": "\nNuclear Power.", "size": 44, "bold": True,
             "color": WHITE, "font": "Georgia"},
        ],
        line_spacing=1.05,
    )

    rect(s, 0.35, 3.55, 6.8, 0.8, TEAL)
    txt(s, 0.5, 3.6, 6.5, 0.7,
        "The Big Idea: Zambia's chronic power deficit is strangling "
        "economic growth — nuclear power is the reliable, scalable, "
        "low-carbon solution that hydro and solar alone cannot provide.",
        size=11.5, italic=True, color=WHITE, line_spacing=1.25)

    txt(s, 0.35, 4.95, 9.0, 0.35,
        "A Case Built on Evidence · Economics · Science · Policy",
        size=10, color=SLATE)


# ---------- SLIDE 2 — Act I: The Problem ----------
def slide_02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, BG_BLUE)
    header_bar(s, NAVY,
               "ACT I — THE PROBLEM", GOLD,
               "Zambia Is Running on Empty")

    stats = [
        (0.3, CRISIS, "image-2-1.png", "8–16 hrs",
         "Daily Load-Shedding",
         "Households & industry cut off daily"),
        (3.55, ORANGE, "image-2-2.png", "43%",
         "Electricity Access Rate",
         "Among the lowest in Southern Africa"),
        (6.8, TEAL, "image-2-3.png", "7%",
         "GDP Growth Target",
         "Impossible without reliable power"),
    ]
    for left, accent, icon, big, label, sub in stats:
        rect(s, left, 1.3, 3.05, 3.2, WHITE)
        rect(s, left, 1.3, 3.05, 0.12, accent)
        picture(s, icon, left + 1.15, 1.55, 0.7, 0.7)
        txt(s, left + 0.1, 2.35, 2.85, 0.75,
            big, size=34, bold=True, color=accent, font="Georgia")
        txt(s, left + 0.1, 3.1, 2.85, 0.45,
            label, size=13, bold=True, color=NAVY_DEEP)
        txt(s, left + 0.1, 3.55, 2.85, 0.6,
            sub, size=10.5, color=SLATE)

    rect(s, 0.3, 4.65, 9.4, 0.65, NAVY_DEEP)
    txt(s, 0.5, 4.72, 9.0, 0.5,
        "Over 99% of Zambia's generation capacity is hydropower — "
        "and climate change is drying up the Kariba Dam year after year.",
        size=12, color=WHITE, line_spacing=1.25)


# ---------- SLIDE 3 — Alternatives Comparison ----------
def slide_03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, WHITE)
    header_bar(s, TEAL,
               "WHY ALTERNATIVES FALL SHORT", WHITE,
               "The Alternatives: Necessary but Not Sufficient")

    cols = [
        # (left, fill, title_color, icon, name,
        #   strengths, limitations,
        #   strength_color, limit_color, text_color, body_color)
        (0.2, BG_BLUE, NAVY_DEEP, "image-3-1.png", "Hydro",
         ["Already installed (2,800 MW)", "Low operating cost"],
         ["Drought-dependent: Kariba at 12% capacity (2024)",
          "Near maximum expansion potential",
          "Climate risk worsening"],
         TEAL, CRISIS, NAVY_DEEP, BODY_GREY),
        (2.65, BG_BLUE, NAVY_DEEP, "image-3-2.png", "Solar",
         ["High resource potential", "Falling panel costs"],
         ["Intermittent — no sun at night",
          "Requires expensive storage at scale",
          "Cannot meet industrial baseload"],
         TEAL, CRISIS, NAVY_DEEP, BODY_GREY),
        (5.1, BG_BLUE, NAVY_DEEP, "image-3-3.png", "Coal / Gas",
         ["Reliable baseload", "Short build time"],
         ["High CO₂ emissions",
          "Climate financing blocked",
          "Fuel import dependency"],
         TEAL, CRISIS, NAVY_DEEP, BODY_GREY),
        (7.55, TEAL, WHITE, "image-3-4.png", "Nuclear",
         ["24/7 baseload — weather-independent",
          "Lowest lifecycle carbon (12g CO₂/kWh)",
          "Smallest land footprint per kWh"],
         ["High upfront capital",
          "Long build time (10–15 yrs)",
          "Requires regulatory capacity"],
         TEAL_MINT, PINK_LIGHT, WHITE, PINK_PALE),
    ]

    for (left, fill, title_color, icon, name, strengths, limits,
         s_color, l_color, t_color, body_color) in cols:
        rect(s, left, 1.25, 2.35, 4.0, fill)
        picture(s, icon, left + 0.87, 1.35, 0.6, 0.6)
        txt(s, left + 0.1, 2.0, 2.15, 0.45,
            name, size=16, bold=True, color=title_color, font="Georgia")
        txt(s, left + 0.1, 2.5, 2.15, 0.3,
            "✓ Strengths", size=9.5, bold=True, color=s_color)
        y = 2.82
        for line in strengths:
            txt(s, left + 0.15, y, 2.05, 0.35,
                line, size=9, color=t_color, line_spacing=1.15)
            y += 0.38
        # limitations heading position
        txt(s, left + 0.1, 3.62, 2.15, 0.3,
            "✗ Limitations", size=9.5, bold=True, color=l_color)
        y = 3.95
        for line in limits:
            txt(s, left + 0.15, y, 2.05, 0.3,
                line, size=8.5, color=body_color, line_spacing=1.15)
            y += 0.32


# ---------- SLIDE 4 — Act II: The Economics ----------
def slide_04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, NAVY_DEEP)

    rect(s, 0, 0, 10, 1.1, GOLD)
    txt(s, 0.5, 0.08, 9.0, 0.35,
        "ACT II — THE ECONOMICS", size=9.5, bold=True, color=NAVY_DEEP)
    txt(s, 0.5, 0.42, 9.0, 0.55,
        "The Cost of Doing Nothing Is Greater",
        size=28, bold=True, color=NAVY_DEEP, font="Georgia")

    # LCOE chart
    data = CategoryChartData()
    data.categories = ["Coal", "Gas", "Nuclear",
                       "Solar+Storage", "Hydro (new)"]
    data.add_series("LCOE (USD/MWh)", (110, 95, 80, 145, 90))
    chart_shape = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(1.25), Inches(5.5), Inches(3.5), data
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Levelised Cost of Electricity (USD/MWh)"
    for p in chart.chart_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.name = "Calibri"
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(10)
    dl.font.bold = True
    dl.font.color.rgb = WHITE
    color_points(chart.series[0],
                 [SLATE, SLATE, TEAL_GREEN, SLATE, SLATE])
    for axis in (chart.category_axis, chart.value_axis):
        try:
            tl = axis.tick_labels
            tl.font.size = Pt(9)
            tl.font.color.rgb = WHITE
        except Exception:
            pass

    # Right-hand cards
    cards = [
        (1.25, "image-4-1.png", "$2B+ annual GDP loss", GOLD,
         "from load-shedding — Zambia's economy loses more every "
         "year than a nuclear plant costs to operate."),
        (2.40, "image-4-2.png", "Green financing unlocked", TEAL_GREEN,
         "Nuclear qualifies for climate financing, SDG bonds & "
         "development bank funding unavailable for coal."),
        (3.55, "image-4-3.png", "50–60 year asset life", TEAL,
         "Once built, nuclear provides generations of cheap, stable "
         "baseload power at $20–30/MWh operational cost."),
    ]
    for top, icon, head, hcolor, body in cards:
        rect(s, 6.1, top, 3.65, 1.0, NAVY)
        picture(s, icon, 6.2, top + 0.25, 0.45, 0.45)
        txt(s, 6.75, top + 0.08, 2.9, 0.35,
            head, size=12, bold=True, color=hcolor)
        txt(s, 6.75, top + 0.46, 2.9, 0.5,
            body, size=9.5, color=SLATE, line_spacing=1.25)

    rect(s, 0, 4.95, 10, 0.65, TEAL)
    txt(s, 0.3, 5.02, 9.4, 0.52,
        "Morocco, Egypt, South Africa & Ghana are all actively pursuing "
        "nuclear — Zambia risks being left behind in the energy "
        "transition.",
        size=11.5, color=WHITE, line_spacing=1.25)


# ---------- SLIDE 5 — Safety & Capability ----------
def slide_05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, BG_BLUE)
    header_bar(s, NAVY, "SAFETY & CAPABILITY", GOLD,
               "Is Zambia Ready? Yes — With the Right Foundation.")

    # LEFT — foundations
    rect(s, 0.3, 1.25, 4.35, 4.0, WHITE)
    txt(s, 0.5, 1.35, 3.95, 0.45,
        "Zambia's Existing Foundations",
        size=14, bold=True, color=NAVY)

    foundations = [
        ("image-5-1.png", "IAEA member state since 1965"),
        ("image-5-2.png",
         "Zambia Centre for Energy Research (ZCER) operational"),
        ("image-5-3.png",
         "Uranium deposits — Zambia has the raw material"),
        ("image-5-4.png", "Active nuclear medicine programme (ZMC)"),
        ("image-5-5.png",
         "Bilateral agreements with Russia (Rosatom) & China"),
        ("image-5-6.png",
         "Engineers trained at SA NNR and regional programs"),
    ]
    y = 1.9
    for icon, text in foundations:
        picture(s, icon, 0.5, y + 0.05, 0.3, 0.3)
        txt(s, 0.9, y, 3.6, 0.42,
            text, size=11, color=NAVY_DEEP)
        y += 0.44

    # RIGHT — SMRs
    rect(s, 5.0, 1.25, 4.65, 4.0, TEAL)
    txt(s, 5.2, 1.35, 4.25, 0.65,
        "Small Modular Reactors (SMRs): The Right-Size Solution",
        size=14, bold=True, color=WHITE)
    txt(s, 5.2, 2.1, 4.25, 0.65,
        "Zambia doesn't need a 1,000 MW Chernobyl-era reactor. "
        "Modern SMRs are purpose-built for developing nations:",
        size=10.5, color=TEAL_LIGHT, line_spacing=1.25)

    smr_lines = [
        "→  50–300 MW output — scalable to Zambia's grid",
        "→  Factory-built, modular — lower build risk",
        "→  Passive safety systems — no meltdown without power",
        "→  Generation IV designs produce 80% less waste",
        ("→  Countries like Jordan, Ghana, and Kenya are already in "
         "advanced planning stages"),
    ]
    y = 2.85
    for line in smr_lines:
        txt(s, 5.2, y, 4.2, 0.42,
            line, size=10.5, color=WHITE, line_spacing=1.2)
        y += 0.44


# ---------- SLIDE 6 — Honest Challenges ----------
def slide_06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, WHITE)
    header_bar(s, CRISIS, "HONEST CHALLENGES — AND SOLUTIONS", WHITE,
               "The Risks Are Real — and Manageable")

    cards = [
        (0.25, 1.25, CRISIS, "image-6-1.png", "High Capital Cost",
         "⚠  A 300 MW SMR costs $1–2B upfront — "
         "significant for Zambia's budget.",
         "✓  Phased financing via World Bank, AIIB, AfDB, and "
         "bilateral loans (Russia, China, UAE). Public-private "
         "partnerships spread the risk."),
        (0.25, 2.67, ORANGE, "image-6-2.png", "Skills Shortage",
         "⚠  Nuclear engineering requires specialist expertise "
         "Zambia currently lacks at scale.",
         "✓  Begin NOW: scholarship programs, partner with IAEA's "
         "nuclear training centres, UNZA nuclear curriculum. Build "
         "the pipeline over 10 years."),
        (0.25, 4.09, NAVY, "image-6-3.png", "Waste Management",
         "⚠  Radioactive waste must be stored safely for centuries.",
         "✓  SMRs produce far less waste. Dry cask storage is "
         "proven. IAEA provides full technical guidance. Zambia has "
         "suitable geological sites."),
        (5.1, 1.25, GOLD, "image-6-4.png", "Safety & Public Perception",
         "⚠  Public fear (Chernobyl, Fukushima) and political "
         "opposition are real barriers.",
         "✓  Independent regulator (like SA's NNR), full "
         "transparency, and data: nuclear is statistically safer per "
         "kWh than coal, gas, or even rooftop solar."),
        (5.1, 2.67, TEAL, "image-6-5.png", "Long Build Timeline",
         "⚠  10–15 years from decision to first power — "
         "a long time to wait.",
         "✓  Start today. The grid won't fix itself. Every year of "
         "delay is another decade of load-shedding. Parallel "
         "solar/wind investments fill the gap."),
        (5.1, 4.09, TEAL_GREEN, "image-6-6.png", "Regulatory Gaps",
         "⚠  Zambia lacks a standalone nuclear regulatory authority.",
         "✓  Establish the Zambia Nuclear Regulatory Authority "
         "(ZaNRA) now. Model on Kenya's KNEB or SA's NNR. "
         "IAEA will assist."),
    ]
    for left, top, accent, icon, head, warn, fix in cards:
        rect(s, left, top, 4.6, 1.3, BG_BLUE)
        rect(s, left, top, 0.1, 1.3, accent)
        picture(s, icon, left + 0.20, top + 0.50, 0.35, 0.35)
        txt(s, left + 0.65, top + 0.08, 3.8, 0.32,
            head, size=12.5, bold=True, color=NAVY_DEEP)
        txt(s, left + 0.65, top + 0.42, 3.8, 0.38,
            warn, size=9, color=CRISIS, line_spacing=1.20)
        txt(s, left + 0.65, top + 0.82, 3.8, 0.38,
            fix, size=9, color=TEAL, line_spacing=1.20)


# ---------- SLIDE 7 — National Development ----------
def slide_07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, NAVY)
    header_bar(s, TEAL, "ACT II — NATIONAL DEVELOPMENT", WHITE,
               "Nuclear Powers More Than the Grid")

    cards = [
        (0.25, 1.3, TEAL_GREEN, "image-7-1.png",
         "Vision 2030 & 8NDP",
         "The 8th National Development Plan calls for "
         "industrialisation and economic transformation. This is "
         "impossible without reliable, affordable power at scale."),
        (5.15, 1.3, GOLD, "image-7-2.png",
         "Climate Commitments",
         "Zambia's NDC pledges low-carbon growth. "
         "Nuclear = 12g CO₂/kWh vs coal at 820g. Meeting climate "
         "goals AND economic growth requires nuclear."),
        (0.25, 3.35, TEAL, "image-7-3.png",
         "Mining & Industrialisation",
         "Copper mining, fertiliser production, and manufacturing "
         "all need 24/7 baseload power. Nuclear anchors Zambia's "
         "industrial future."),
        (5.15, 3.35, ORANGE, "image-7-4.png",
         "STEM & Human Capital",
         "A nuclear programme creates 10,000+ skilled jobs, "
         "elevates Zambia's universities, and builds a generation "
         "of scientists and engineers."),
    ]
    for left, top, accent, icon, head, body in cards:
        rect(s, left, top, 4.55, 1.75, NAVY_PANEL)
        rect(s, left, top, 4.55, 0.1, accent)
        picture(s, icon, left + 0.2, top + 0.3, 0.5, 0.5)
        txt(s, left + 0.85, top + 0.18, 3.55, 0.4,
            head, size=13, bold=True, color=accent)
        txt(s, left + 0.85, top + 0.62, 3.55, 1.0,
            body, size=10, color=SLATE, line_spacing=1.25)


# ---------- SLIDE 8 — Roadmap (3 phases) ----------
def slide_08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, BG_BLUE)
    header_bar(s, NAVY, "ACT III — THE SOLUTION", GOLD,
               "A 3-Phase Roadmap to Nuclear Power")

    phases = [
        (0.25, TEAL, "PHASE 1", "2025–2030", "Foundations",
         ["Establish ZaNRA (Regulatory Authority)",
          "Join IAEA milestone process",
          "Launch nuclear engineering scholarships",
          "Conduct site surveys & feasibility studies",
          "Build public awareness campaign"]),
        (3.5, GOLD, "PHASE 2", "2030–2038", "Construction",
         ["Select SMR technology & partner",
          "Secure financing package",
          "Begin construction of first 300 MW SMR",
          "Build local supply chain",
          "Train 500+ nuclear technicians"]),
        (6.75, TEAL_GREEN, "PHASE 3", "2038–2050+",
         "Operation & Expansion",
         ["First power to the grid (~2038)",
          "Expand to 1,000+ MW by 2045",
          "Export surplus power regionally via SAPP",
          "Fund second plant from revenues",
          "Become a regional nuclear hub"]),
    ]
    for left, accent, label, dates, title, items in phases:
        rect(s, left, 1.3, 3.1, 4.0, WHITE)
        rect(s, left, 1.3, 3.1, 0.85, accent)
        txt(s, left + 0.15, 1.35, 2.8, 0.35,
            label, size=11, bold=True, color=WHITE)
        txt(s, left + 0.15, 1.72, 2.8, 0.3,
            dates, size=10.5, color=TEAL_PALE)
        txt(s, left + 0.15, 2.2, 2.8, 0.4,
            title, size=16, bold=True, color=accent, font="Georgia")
        y = 2.68
        for line in items:
            textbox(
                s, left + 0.15, y, 2.8, 0.46,
                [
                    {"text": "→  ", "size": 10, "bold": True,
                     "color": accent},
                    {"text": line, "size": 10, "color": NAVY_DEEP},
                ],
                line_spacing=1.15,
            )
            y += 0.48


# ---------- SLIDE 9 — Big Idea + Recommendations ----------
def slide_09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, NAVY_DEEP)

    picture(s, "image-9-1.png", 5.5, 0.2, 4.5, 4.5)
    rect(s, 0, 0, 0.12, 5.625, GOLD)

    txt(s, 0.35, 0.5, 7.0, 0.4,
        "THE BIG IDEA", size=11, bold=True, color=GOLD)

    txt(s, 0.35, 1.0, 8.8, 2.5,
        "\"Zambia cannot achieve its Vision 2030 development goals on "
        "a grid powered by rain alone — nuclear power is the decisive "
        "investment that ends load-shedding, anchors "
        "industrialisation, and positions Zambia as Africa's next "
        "energy leader.\"",
        size=20, italic=True, color=WHITE, font="Georgia",
        line_spacing=1.3)

    asks = [
        (3.55, GOLD, "01",
         "Commission a national nuclear feasibility study NOW"),
        (4.13, TEAL, "02",
         "Establish the Zambia Nuclear Regulatory Authority"),
        (4.71, TEAL_GREEN, "03",
         "Begin IAEA milestone consultations within 12 months"),
    ]
    for top, color, num, ask in asks:
        rect(s, 0.35, top, 0.5, 0.45, color)
        txt(s, 0.35, top + 0.05, 0.5, 0.35,
            num, size=12, bold=True, color=NAVY_DEEP,
            align=PP_ALIGN.CENTER)
        txt(s, 0.97, top + 0.06, 8.0, 0.38,
            ask, size=13, bold=True, color=WHITE)

    txt(s, 0.35, 5.22, 9.0, 0.3,
        "The lights don't have to keep going off.",
        size=10.5, italic=True, color=SLATE)


# ---------- SLIDE 10 — Sources ----------
def slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    set_bg(s, BG_BLUE)

    rect(s, 0, 0, 10, 0.9, NAVY)
    txt(s, 0.5, 0.27, 9.0, 0.4,
        "KEY EVIDENCE SOURCES & DATA NOTES",
        size=14, bold=True, color=WHITE)

    sources = [
        "1.  International Atomic Energy Agency (IAEA) — Nuclear "
        "Energy Basic Principles & SMR Technology Catalogue 2023",
        "2.  World Nuclear Association — Nuclear Power Economics "
        "2024; LCOE comparisons",
        "3.  Zambia Electricity Supply Corporation (ZESCO) — Annual "
        "Reports 2022–2024; Generation capacity statistics",
        "4.  International Energy Agency (IEA) — Africa Energy "
        "Outlook 2023; Zambia grid access data (43%)",
        "5.  World Bank — Zambia Energy Sector Assessment 2023; "
        "Load-shedding economic impact ($2B+ GDP)",
        "6.  IPCC AR6 — Life-cycle GHG emissions by technology: "
        "Nuclear 12 gCO₂eq/kWh vs. Coal 820g",
        "7.  Zambia Ministry of Energy — 8th National Development "
        "Plan (8NDP) energy targets; Vision 2030",
        "8.  IAEA — Country Nuclear Power Profiles: Ghana, Kenya, "
        "Egypt — SMR development timelines",
        "9.  Rosatom International — Small Modular Reactor technical "
        "specifications & cost estimates",
        "10.  African Development Bank — Renewable energy financing "
        "and NDC green bond frameworks, 2024",
    ]
    runs = []
    for i, line in enumerate(sources):
        prefix = "\n" if i > 0 else ""
        runs.append({"text": prefix + line, "size": 9.5,
                     "color": NAVY_DEEP})
    textbox(s, 0.4, 1.05, 9.2, 4.2, runs, line_spacing=1.35)

    txt(s, 0.4, 5.18, 9.2, 0.3,
        "Data points where noted are approximations based on published "
        "ranges. Specific Zambia statistics sourced from ZESCO and "
        "Ministry of Energy publications.",
        size=8.5, italic=True, color=SLATE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)

    prs.save(OUTPUT)
    print(f"Generated: {OUTPUT}")


if __name__ == "__main__":
    build()
