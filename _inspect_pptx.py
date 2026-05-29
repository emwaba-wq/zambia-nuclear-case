"""Dump structure of the current .pptx so we can reconcile build_presentation.py."""
from pptx import Presentation
from pptx.util import Emu

P = "The_Case_for_Nuclear_Power_in_Zambia.pptx"


def emu_to_in(v):
    return round(v / 914400, 3) if v is not None else None


def rgb(c):
    try:
        return str(c.rgb) if c and c.type is not None else None
    except Exception:
        return None


prs = Presentation(P)
print(f"# Slide size: "
      f"{emu_to_in(prs.slide_width)} x {emu_to_in(prs.slide_height)} in")
print(f"# Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"=== SLIDE {i} ===")
    try:
        bg = slide.background.fill
        if bg.type is not None:
            print(f"  bg.type={bg.type}", end="")
            try:
                print(f" rgb={rgb(bg.fore_color)}")
            except Exception:
                print()
    except Exception:
        pass
    for sh in slide.shapes:
        kind = sh.shape_type
        L = emu_to_in(sh.left)
        T = emu_to_in(sh.top)
        W = emu_to_in(sh.width)
        H = emu_to_in(sh.height)
        line = f"  [{kind}] L={L} T={T} W={W} H={H}"
        try:
            if sh.has_text_frame:
                txt = " | ".join(
                    p.text for p in sh.text_frame.paragraphs if p.text
                )
                if txt:
                    snippet = txt[:160].replace("\n", " ")
                    line += f"  TXT='{snippet}'"
                # Capture first run font info
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        f = r.font
                        line += (
                            f"\n      run: '{r.text[:60]}'"
                            f" name={f.name} size={f.size}"
                            f" bold={f.bold} italic={f.italic}"
                            f" color={rgb(f.color) if f.color else None}"
                        )
                        break
                    break
        except Exception as e:
            line += f"  (text-error: {e})"
        try:
            if hasattr(sh, "fill") and sh.fill.type is not None:
                fc = rgb(sh.fill.fore_color)
                if fc:
                    line += f"  FILL={fc}"
        except Exception:
            pass
        print(line)
    # Notes
    try:
        ntf = slide.notes_slide.notes_text_frame
        ntxt = "\n".join(p.text for p in ntf.paragraphs if p.text)
        if ntxt.strip():
            print(f"  NOTES: {ntxt[:500]}")
    except Exception:
        pass
    print()
